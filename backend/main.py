"""
Negosyo Plan – FastAPI backend
- DeepSeek AI blueprint generation
- Admin authentication (JWT) + site settings CRUD

Run:
    uvicorn main:app --reload --port 8000
"""

import hmac
import json
import os
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Optional

import shutil
import uuid

import httpx
import jwt as pyjwt
from dotenv import load_dotenv
from fastapi import FastAPI, File, Header, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from fastapi.responses import FileResponse

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

load_dotenv()

# ── Config ────────────────────────────────────────────────────────────────────

DEEPSEEK_API_KEY: str = os.getenv("DEEPSEEK_API_KEY", "")
DEEPSEEK_URL = "https://api.deepseek.com/chat/completions"
DEEPSEEK_MODEL = "deepseek-chat"

GENSPARK_API_KEY: str = os.getenv("GENSPARK_API_KEY", "")
GENSPARK_URL = "https://api.genspark.ai/v1/agents/sparks"

# Admin auth
ADMIN_USERNAME: str = os.getenv("ADMIN_USERNAME", "admin")
ADMIN_PASSWORD: str = os.getenv("ADMIN_PASSWORD", "NegosyoPlan@2026!")
JWT_SECRET: str = os.getenv("JWT_SECRET", "np-jwt-secret-change-in-production-2026")
JWT_ALGORITHM = "HS256"
JWT_EXPIRE_HOURS = 10

# Settings file (persisted on disk next to main.py)
SETTINGS_PATH = Path(__file__).parent / "settings.json"
ASSETS_IMAGES = Path(__file__).parent.parent / "assets" / "images"
PRESENTATIONS_DIR = Path(__file__).parent.parent / "generated_presentations"

DEFAULT_SETTINGS: dict[str, Any] = {
    "pricing": {
        "starter": 1900,
        "founder": 2600,
        "complete": 3500,
        "custom": 4600,
    },
    "promo_banner": {
        "enabled": False,
        "text": "🎉 Special launch offer — 20% off all bundles this week!",
        "bg_color": "#E8420A",
        "text_color": "#ffffff",
        "link": "shop.html",
        "link_text": "Shop Now",
    },
    "theme": {
        "primary": "#E8420A",
        "secondary": "#0F1F3D",
        "accent": "#F5A500",
    },
    "hero": {
        "title": "Launch your business with a ready-made blueprint.",
        "subtitle": "Digital business blueprints for OFWs, entrepreneurs, and startup creators.",
        "cta_primary": "Browse Bundles",
        "cta_secondary": "Generate AI Blueprint",
    },
    "products": {
        "1": {
            "name": "Starter Bundle",
            "image": "https://images.unsplash.com/photo-1554224155-6726b3ff858f?ixlib=rb-4.0.3&auto=format&fit=crop&w=400&q=80",
            "description": "Perfect for beginners starting their entrepreneurial journey.",
        },
        "2": {
            "name": "Founder Bundle",
            "image": "https://images.unsplash.com/photo-1519389950473-47ba0277781c?ixlib=rb-4.0.3&auto=format&fit=crop&w=400&q=80",
            "description": "Advanced business strategy and investor-ready templates.",
        },
        "3": {
            "name": "Complete Set Bundle",
            "image": "https://images.unsplash.com/photo-1554224155-8d04cb21cd6c?ixlib=rb-4.0.3&auto=format&fit=crop&w=400&q=80",
            "description": "Comprehensive package for planning, launch, and growth.",
        },
        "4": {
            "name": "Custom Bundle",
            "image": "https://images.unsplash.com/photo-1556761175-4b46a572b786?ixlib=rb-4.0.3&auto=format&fit=crop&w=400&q=80",
            "description": "Tailor-made toolkit for your market, operations, and goals.",
        },
    },
    "contact": {
        "whatsapp": "1234567890",
        "email": "support@negosyoplan.com",
    },
    "branding": {
        "logo_url": "",
        "brand_name": "NEGOSYO PLAN",
    },
}


# ── Settings helpers ──────────────────────────────────────────────────────────

def load_settings() -> dict:
    if SETTINGS_PATH.exists():
        try:
            return json.loads(SETTINGS_PATH.read_text(encoding="utf-8"))
        except Exception:
            pass
    return json.loads(json.dumps(DEFAULT_SETTINGS))   # deep copy


def save_settings(data: dict) -> None:
    SETTINGS_PATH.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def deep_merge(base: dict, override: dict) -> dict:
    result = base.copy()
    for k, v in override.items():
        if k in result and isinstance(result[k], dict) and isinstance(v, dict):
            result[k] = deep_merge(result[k], v)
        else:
            result[k] = v
    return result


# ── JWT helpers ───────────────────────────────────────────────────────────────

def create_token(username: str) -> str:
    exp = datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRE_HOURS)
    return pyjwt.encode({"sub": username, "exp": exp}, JWT_SECRET, algorithm=JWT_ALGORITHM)


def decode_token(token: str) -> str | None:
    try:
        payload = pyjwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        return payload.get("sub")
    except pyjwt.PyJWTError:
        return None


def require_admin(authorization: str | None) -> str:
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(403, "Not authenticated")
    user = decode_token(authorization.split(" ", 1)[1])
    if not user:
        raise HTTPException(401, "Token expired or invalid — please log in again")
    return user

# ── App ───────────────────────────────────────────────────────────────────────

app = FastAPI(
    title="Negosyo Plan AI API",
    description="DeepSeek-powered business blueprint generator for Negosyo Plan",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],          # tighten to your domain in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Request / Response models ─────────────────────────────────────────────────

class LocationData(BaseModel):
    country: str = "Philippines"
    region: str = ""
    city: str = ""
    subdivision: str = ""


class SupplierInfo(BaseModel):
    name: str
    type: str = ""
    contact: str = ""
    coverage: str = ""


class BlueprintRequest(BaseModel):
    business_type: str = Field(..., description="Slug e.g. mini_grocery, restaurant")
    business_type_label: str = Field(..., description="Human-readable label")
    bundle_tier: str = Field(
        "starter", description="starter | founder | complete | custom"
    )
    capital: float = Field(..., description="Available capital in PHP")
    location: LocationData = Field(default_factory=LocationData)
    description: Optional[str] = Field("", description="Buyer's own description")
    suppliers: list[SupplierInfo] = Field(default_factory=list)


class QuickAnalysisRequest(BaseModel):
    business_type: str
    business_type_label: str
    bundle_tier: str = "starter"
    capital: float
    location: LocationData = Field(default_factory=LocationData)


# ── Prompt helpers ────────────────────────────────────────────────────────────

TIER_DEPTH = {
    "starter": (
        "foundational / essential level — practical basics, quick-start checklists, "
        "beginner-friendly language. Skip advanced frameworks."
    ),
    "founder": (
        "advanced level — include Porter's Five Forces, investor-ready financial "
        "analysis, competitive positioning and market validation."
    ),
    "complete": (
        "professional / pro level — full depth including PESTEL, ANSOFF matrix, "
        "detailed 12-month financial modelling and executive-grade strategy."
    ),
    "custom": (
        "premium custom level — maximum depth, highly personalised to the exact "
        "business, location, capital, and selected suppliers."
    ),
}

SYSTEM_PROMPT = """You are a senior Filipino business consultant with 20+ years' experience advising
OFW (Overseas Filipino Worker) entrepreneurs and MSMEs across all Philippine regions.

Your expertise covers:
• DTI/BIR/SEC registration, barangay permits, local government licensing
• Regional supplier ecosystems (NCR, Visayas, Mindanao, Luzon provinces)
• SB Finance, LANDBANK, DBP, Pag-IBIG MSME loan programs
• Filipino consumer behaviour, seasonal demand, fiestas, remittance cycles
• Both traditional (tarpaulin, text brigades, barangay events) and digital marketing
  (Facebook Shops, TikTok Live, Shopee, Lazada, GCash QR payments)

IMPORTANT: Respond with valid JSON only. No markdown fences, no prose outside the JSON."""


def _location_str(loc: LocationData) -> str:
    parts = [loc.subdivision, loc.city, loc.region, loc.country]
    return ", ".join(p for p in parts if p)


def _supplier_names(suppliers: list[SupplierInfo]) -> str:
    if not suppliers:
        return "none specified"
    return ", ".join(s.name for s in suppliers)


# ── Full blueprint prompt ─────────────────────────────────────────────────────

def build_blueprint_prompt(req: BlueprintRequest) -> str:
    depth = TIER_DEPTH.get(req.bundle_tier, TIER_DEPTH["starter"])
    loc = _location_str(req.location)
    cap = f"PHP {req.capital:,.0f}"
    sups = _supplier_names(req.suppliers)

    is_founder_plus = req.bundle_tier in ("founder", "complete", "custom")
    is_complete_plus = req.bundle_tier in ("complete", "custom")

    porter = (
        '"porters_five_forces": {'
        '"supplier_power": "string",'
        '"buyer_power": "string",'
        '"competitive_rivalry": "string",'
        '"threat_of_substitutes": "string",'
        '"threat_of_new_entrants": "string",'
        '"overall_assessment": "string"'
        '},'
    ) if is_founder_plus else ""

    pestel_ansoff = (
        '"pestel_analysis": {'
        '"political": "string","economic": "string","social": "string",'
        '"technological": "string","environmental": "string","legal": "string"'
        '},'
        '"ansoff_matrix": {'
        '"market_penetration": "string","market_development": "string",'
        '"product_development": "string","diversification": "string"'
        '},'
    ) if is_complete_plus else ""

    return f"""Analyse this business opportunity and return a comprehensive JSON blueprint at {depth}.

INPUTS
  Business type : {req.business_type_label} ({req.business_type})
  Bundle tier   : {req.bundle_tier.title()}
  Capital       : {cap}
  Location      : {loc}
  Description   : {req.description or "Not provided"}
  Suppliers     : {sups}

Return ONLY the following JSON structure. Fill every "string" with specific,
actionable, location-aware content. Fill every number with a realistic PHP figure.

{{
  "executive_summary": {{
    "business_concept": "string",
    "unique_value_proposition": "string",
    "target_market": "string",
    "revenue_model": "string",
    "growth_potential": "string",
    "key_success_factors": ["string","string","string"]
  }},
  "swot_analysis": {{
    "strengths":    ["string","string","string","string","string"],
    "weaknesses":   ["string","string","string","string"],
    "opportunities":["string","string","string","string","string"],
    "threats":      ["string","string","string","string"],
    "so_strategies":["string","string","string"],
    "wo_strategies":["string","string"],
    "st_strategies":["string","string"],
    "wt_strategies":["string","string"]
  }},
  "financial_analysis": {{
    "startup_cost_breakdown": [
      {{"item":"string","estimated_cost_php":0,"notes":"string"}}
    ],
    "monthly_operating_costs": [
      {{"item":"string","estimated_cost_php":0,"notes":"string"}}
    ],
    "revenue_projections": {{
      "month_1":0,"month_2":0,"month_3":0,
      "month_4":0,"month_5":0,"month_6":0
    }},
    "break_even_months": 0,
    "gross_margin_percent": 0,
    "roi_12_months_percent": 0,
    "financing_options": ["string","string","string"]
  }},
  "management_system": {{
    "daily_operations_checklist": [
      "string","string","string","string","string","string","string","string"
    ],
    "key_performance_indicators": [
      {{"kpi":"string","target":"string","frequency":"string"}}
    ],
    "staffing_plan": [
      {{"role":"string","count":0,"monthly_salary_php":0,"responsibilities":"string"}}
    ],
    "standard_operating_procedures": ["string","string","string","string"]
  }},
  "fifo_inventory": {{
    "applicable": true,
    "product_categories": [
      {{"category":"string","rotation_rule":"string","shelf_life":"string","reorder_point":"string"}}
    ],
    "stock_management_tips": ["string","string","string","string"],
    "abc_classification": {{"a_items":"string","b_items":"string","c_items":"string"}}
  }},
  "risk_analysis": {{
    "top_risks": [
      {{"risk":"string","probability":"High|Medium|Low","impact":"High|Medium|Low","mitigation":"string"}}
    ],
    "contingency_plans": ["string","string","string"],
    "insurance_recommendations": ["string","string"]
  }},
  "business_strategy": {{
    "competitive_advantage": "string",
    "pricing_strategy": "string",
    "value_proposition_canvas": {{
      "customer_jobs":    ["string","string","string"],
      "customer_pains":   ["string","string","string"],
      "customer_gains":   ["string","string","string"],
      "pain_relievers":   ["string","string","string"],
      "gain_creators":    ["string","string","string"],
      "products_services":["string","string","string"]
    }},
    {porter}
    "growth_milestones": [
      {{"month":0,"milestone":"string","action":"string"}}
    ]
  }},
  {pestel_ansoff}
  "traditional_marketing": {{
    "strategies": [
      {{"tactic":"string","description":"string","monthly_budget_php":0,"expected_reach":"string"}}
    ],
    "local_partnerships": ["string","string","string"],
    "offline_channels":   ["string","string","string"]
  }},
  "digital_marketing": {{
    "platforms": [
      {{"platform":"string","strategy":"string","content_type":"string",
        "posting_frequency":"string","monthly_budget_php":0}}
    ],
    "seo_keywords":     ["string","string","string","string","string"],
    "email_sms_strategy":"string",
    "online_marketplaces":["string","string"]
  }},
  "implementation_timeline": [
    {{"phase":0,"title":"string","duration":"string",
      "key_tasks":["string","string","string"],"budget_php":0}}
  ],
  "supplier_recommendations": [
    {{"name":"string","type":"string","contact":"string",
      "coverage":"string","why_recommended":"string"}}
  ],
  "legal_requirements": {{
    "registrations": ["string","string","string","string"],
    "permits":        ["string","string","string"],
    "compliance":     ["string","string"]
  }},
  "permit_acquisition_guide": {{
    "steps": [
      {{
        "step": 1,
        "requirement": "string",
        "office": "string",
        "documents_needed": ["string","string","string"],
        "cost_php_low": 0,
        "cost_php_high": 0,
        "processing_days": "string",
        "where_to_file": "string",
        "tips": "string"
      }}
    ],
    "total_estimated_cost_php_low": 0,
    "total_estimated_cost_php_high": 0,
    "total_processing_weeks": "string",
    "important_reminders": ["string","string","string"],
    "online_portals": ["string","string"]
  }},
  "materials_and_equipment": {{
    "equipment": [
      {{
        "name": "string",
        "description": "string",
        "cost_php_low": 0,
        "cost_php_high": 0,
        "priority": "Essential|Important|Nice-to-have",
        "where_to_buy": "string",
        "specs": "string"
      }}
    ],
    "raw_materials": [
      {{
        "name": "string",
        "unit": "string",
        "unit_cost_php": 0,
        "monthly_quantity": 0,
        "monthly_cost_php": 0
      }}
    ],
    "consumables_monthly_php": 0,
    "total_equipment_cost_php_low": 0,
    "total_equipment_cost_php_high": 0,
    "total_monthly_materials_php": 0,
    "sourcing_tips": ["string","string","string"]
  }}
}}"""


# ── Quick analysis prompt ─────────────────────────────────────────────────────

def build_quick_prompt(req: QuickAnalysisRequest) -> str:
    loc = _location_str(req.location)
    cap = f"PHP {req.capital:,.0f}"
    return f"""Quick feasibility check:
  Business: {req.business_type_label} in {loc}
  Capital : {cap}
  Tier    : {req.bundle_tier}

Return ONLY this JSON:
{{
  "feasibility_score": 0,
  "feasibility_label": "Excellent|Good|Moderate|Challenging",
  "quick_summary": "string (2-3 sentences)",
  "top_3_strengths": ["string","string","string"],
  "top_3_risks":     ["string","string","string"],
  "monthly_revenue_estimate_php": 0,
  "break_even_months": 0,
  "recommended_bundle": "starter|founder|complete|custom",
  "bundle_reason": "string",
  "quick_tips": ["string","string","string"],
  "capital_sufficiency": "Sufficient|Tight|Insufficient",
  "capital_note": "string"
}}"""


# ── DeepSeek helper ───────────────────────────────────────────────────────────

async def call_deepseek(
    system: str,
    user: str,
    max_tokens: int = 6000,
    temperature: float = 0.65,
    timeout: float = 90.0,
) -> dict:
    if not DEEPSEEK_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="DEEPSEEK_API_KEY is not set. Add it to backend/.env",
        )

    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        "response_format": {"type": "json_object"},
        "temperature": temperature,
        "max_tokens": max_tokens,
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json",
    }

    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            resp = await client.post(DEEPSEEK_URL, json=payload, headers=headers)
            resp.raise_for_status()
    except httpx.TimeoutException:
        raise HTTPException(504, "DeepSeek request timed out — please retry.")
    except httpx.HTTPStatusError as exc:
        raise HTTPException(
            exc.response.status_code,
            f"DeepSeek API error: {exc.response.text[:300]}",
        )

    raw = resp.json()["choices"][0]["message"]["content"]
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        raise HTTPException(500, "DeepSeek returned non-JSON content.")


# ── Presentation generation ───────────────────────────────────────────────────

_ORANGE = "#E8420A"
_NAVY   = "#0F1F3D"
_AMBER  = "#F5A500"
_WHITE  = "#FFFFFF"
_LIGHT  = "#F4F6F8"
_MUTED  = "#64748B"
_GREEN  = "#16A34A"
_RED    = "#DC2626"
_DARK   = "#1E293B"


def _c(h: str) -> "RGBColor":
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rect(slide, l: float, t: float, w: float, h: float, fill: str, line: bool = False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _c(fill)
    if line:
        shape.line.color.rgb = _c(fill)
    else:
        shape.line.fill.background()
    return shape


def _tb(slide, text: str, l: float, t: float, w: float, h: float,
        sz: int, bold: bool, color: str, align=None, wrap: bool = True):
    if align is None:
        align = PP_ALIGN.LEFT
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = _c(color)
    return txb


def _bullets(slide, items: list, l: float, t: float, w: float, h: float,
             sz: int, color: str, max_n: int = 8):
    items = [str(x) for x in (items or [])[:max_n]]
    if not items:
        return
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(sz)
        run.font.color.rgb = _c(color)


def _header(slide, title: str, sub: str = ""):
    _rect(slide, 0, 0, 13.33, 0.9, _NAVY)
    _rect(slide, 0, 0.9, 13.33, 0.05, _ORANGE)
    _tb(slide, title, 0.3, 0.08, 10, 0.5, 22, True, _WHITE, PP_ALIGN.LEFT)
    if sub:
        _tb(slide, sub, 0.3, 0.55, 12.7, 0.32, 10, False, _AMBER, PP_ALIGN.LEFT)


def generate_pptx(blueprint: dict, meta: dict) -> str:
    PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)

    prs = PPTXPresentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    b = blueprint
    biz_label = meta.get("business_type_label", "Business")
    bundle    = meta.get("bundle_tier", "starter").title()
    cap_php   = meta.get("capital_php", 0)
    cap       = f"PHP {cap_php:,.0f}"
    loc_d     = meta.get("location", {})
    loc       = ", ".join(filter(None, [
        loc_d.get("city", ""), loc_d.get("region", ""), loc_d.get("country", "Philippines")
    ]))

    es     = b.get("executive_summary", {})
    swot   = b.get("swot_analysis", {})
    fin    = b.get("financial_analysis", {})
    mgmt   = b.get("management_system", {})
    risk_d = b.get("risk_analysis", {})
    trad   = b.get("traditional_marketing", {})
    dig    = b.get("digital_marketing", {})
    tl     = b.get("implementation_timeline", [])
    perm   = b.get("permit_acquisition_guide", {})
    mats   = b.get("materials_and_equipment", {})

    # ── Slide 1: Cover ──────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _NAVY)
    _rect(sl, 0, 0, 13.33, 0.12, _ORANGE)
    _rect(sl, 0, 6.8, 13.33, 0.7, _ORANGE)
    _rect(sl, 0.35, 2.3, 0.08, 2.5, _AMBER)
    _tb(sl, biz_label.upper(), 0.6, 2.1, 12, 1.2, 38, True, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "Business Blueprint", 0.6, 3.2, 10, 0.6, 20, False, _AMBER, PP_ALIGN.LEFT)
    _tb(sl, f"Capital: {cap}   |   {loc}", 0.6, 3.85, 10, 0.45, 12, False, _MUTED, PP_ALIGN.LEFT)
    _tb(sl, f"{bundle} Bundle  •  Powered by DeepSeek AI", 0.6, 4.3, 10, 0.4, 11, False, "#8B9DC3", PP_ALIGN.LEFT)
    _tb(sl, "Negosyo Plan", 0.35, 6.87, 6, 0.5, 16, True, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "Your Business Blueprint to Success", 6.5, 6.92, 6.5, 0.4, 10, False, "#FFCDA0", PP_ALIGN.RIGHT)

    # ── Slide 2: Executive Summary ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "Executive Summary", es.get("business_concept", "")[:90])
    cols = [
        ("Business Concept",        es.get("business_concept", "N/A"),        _NAVY),
        ("Unique Value Proposition", es.get("unique_value_proposition", "N/A"), "#1E3A5F"),
        ("Revenue Model",           es.get("revenue_model", "N/A"),            "#162D4A"),
    ]
    for i, (ttl, body, bg) in enumerate(cols):
        xl = 0.3 + i * 4.3
        _rect(sl, xl, 1.1, 4.0, 3.3, bg)
        _tb(sl, ttl, xl + 0.12, 1.18, 3.76, 0.42, 10, True, _AMBER, PP_ALIGN.LEFT)
        _tb(sl, str(body)[:220], xl + 0.12, 1.62, 3.76, 2.6, 8.5, False, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "Key Success Factors", 0.3, 4.6, 6, 0.35, 11, True, _NAVY, PP_ALIGN.LEFT)
    _bullets(sl, es.get("key_success_factors", []), 0.3, 4.95, 6.2, 2.3, 9, _DARK, 5)
    _tb(sl, "Growth Potential", 6.8, 4.6, 6, 0.35, 11, True, _NAVY, PP_ALIGN.LEFT)
    _tb(sl, str(es.get("growth_potential", ""))[:260], 6.8, 4.95, 6.2, 2.3, 9, _DARK, PP_ALIGN.LEFT)

    # ── Slide 3: SWOT Analysis ──────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "SWOT Analysis")
    quads = [
        (0.3,  1.1,  _GREEN, "STRENGTHS",    swot.get("strengths", [])),
        (6.85, 1.1,  _RED,   "WEAKNESSES",   swot.get("weaknesses", [])),
        (0.3,  4.25, _AMBER, "OPPORTUNITIES", swot.get("opportunities", [])),
        (6.85, 4.25, _DARK,  "THREATS",      swot.get("threats", [])),
    ]
    for xl, yt, bg, lbl, items in quads:
        _rect(sl, xl, yt, 6.15, 2.98, bg)
        _tb(sl, lbl, xl + 0.15, yt + 0.1, 5.9, 0.42, 11, True, _WHITE, PP_ALIGN.LEFT)
        _bullets(sl, items, xl + 0.15, yt + 0.52, 5.9, 2.28, 8.5, _WHITE, 4)

    # ── Slide 4: Financial Analysis ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "Financial Analysis", "Revenue Projections & Key Metrics")
    rev = fin.get("revenue_projections", {})
    rev_vals = [rev.get(f"month_{i+1}", 0) for i in range(6)]
    _rect(sl, 0.3, 1.1, 12.7, 0.42, _NAVY)
    _tb(sl, "6-Month Revenue Projection", 0.4, 1.14, 8, 0.3, 10, True, _WHITE, PP_ALIGN.LEFT)
    for i, v in enumerate(rev_vals):
        xl = 0.3 + i * 2.1
        _rect(sl, xl, 1.55, 2.05, 0.72, _LIGHT if i % 2 == 0 else _WHITE)
        _tb(sl, f"Month {i+1}", xl + 0.05, 1.6, 1.95, 0.25, 8, True, _DARK, PP_ALIGN.CENTER)
        _tb(sl, f"₱{v:,.0f}", xl + 0.05, 1.86, 1.95, 0.3, 9, False, _ORANGE, PP_ALIGN.CENTER)
    metrics = [
        ("Break-Even",   f"{fin.get('break_even_months', 0)} months",        _NAVY),
        ("Gross Margin", f"{fin.get('gross_margin_percent', 0)}%",            _ORANGE),
        ("12-Month ROI", f"{fin.get('roi_12_months_percent', 0)}%",           _AMBER),
    ]
    for i, (lbl, val, bg) in enumerate(metrics):
        xl = 0.3 + i * 4.3
        _rect(sl, xl, 2.45, 3.9, 0.9, bg)
        _tb(sl, lbl, xl + 0.1, 2.5, 3.7, 0.3, 9, False, _WHITE, PP_ALIGN.CENTER)
        _tb(sl, val, xl + 0.1, 2.82, 3.7, 0.42, 15, True, _WHITE, PP_ALIGN.CENTER)
    costs = fin.get("startup_cost_breakdown", [])[:6]
    _tb(sl, "Startup Cost Breakdown", 0.3, 3.55, 6, 0.34, 11, True, _NAVY, PP_ALIGN.LEFT)
    for i, cost in enumerate(costs):
        yt = 3.9 + i * 0.43
        _rect(sl, 0.3, yt, 5.8, 0.38, _WHITE if i % 2 == 0 else _LIGHT)
        _tb(sl, str(cost.get("item", ""))[:48], 0.4, yt + 0.06, 4.0, 0.27, 8, False, _DARK, PP_ALIGN.LEFT)
        _tb(sl, f"₱{cost.get('estimated_cost_php', 0):,.0f}", 4.5, yt + 0.06, 1.5, 0.27, 8, True, _ORANGE, PP_ALIGN.RIGHT)
    _tb(sl, "Financing Options", 6.8, 3.55, 6, 0.34, 11, True, _NAVY, PP_ALIGN.LEFT)
    _bullets(sl, fin.get("financing_options", []), 6.8, 3.9, 6.2, 3.4, 9, _DARK, 5)

    # ── Slide 5: Operations & Management ────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "Operations & Management")
    _rect(sl, 0.3, 1.1, 5.8, 0.4, _NAVY)
    _tb(sl, "Daily Operations Checklist", 0.4, 1.14, 5.6, 0.3, 10, True, _WHITE, PP_ALIGN.LEFT)
    _bullets(sl, mgmt.get("daily_operations_checklist", []), 0.3, 1.55, 5.8, 3.1, 8.5, _DARK, 8)
    _rect(sl, 6.5, 1.1, 6.5, 0.4, _NAVY)
    _tb(sl, "Staffing Plan", 6.6, 1.14, 6.3, 0.3, 10, True, _WHITE, PP_ALIGN.LEFT)
    for i, s in enumerate((mgmt.get("staffing_plan", []) or [])[:4]):
        yt = 1.55 + i * 0.76
        _rect(sl, 6.5, yt, 6.5, 0.7, _WHITE if i % 2 == 0 else _LIGHT)
        _tb(sl, f"{s.get('role','N/A')} × {s.get('count',1)}", 6.6, yt + 0.04, 4, 0.28, 9, True, _DARK, PP_ALIGN.LEFT)
        _tb(sl, f"₱{s.get('monthly_salary_php',0):,.0f}/mo", 10.8, yt + 0.04, 2.0, 0.28, 9, False, _ORANGE, PP_ALIGN.RIGHT)
        _tb(sl, str(s.get("responsibilities",""))[:72], 6.6, yt + 0.34, 6.3, 0.28, 7.5, False, _MUTED, PP_ALIGN.LEFT)
    _rect(sl, 0.3, 4.75, 12.7, 0.4, _ORANGE)
    _tb(sl, "Key Performance Indicators", 0.4, 4.79, 8, 0.3, 10, True, _WHITE, PP_ALIGN.LEFT)
    for i, kpi in enumerate((mgmt.get("key_performance_indicators", []) or [])[:4]):
        xl = 0.3 + i * 3.2
        _rect(sl, xl, 5.2, 3.1, 2.1, _WHITE)
        _tb(sl, str(kpi.get("kpi",""))[:38], xl + 0.1, 5.25, 2.9, 0.4, 9, True, _NAVY, PP_ALIGN.LEFT)
        _tb(sl, f"Target: {kpi.get('target','N/A')}", xl + 0.1, 5.65, 2.9, 0.32, 8, False, _ORANGE, PP_ALIGN.LEFT)
        _tb(sl, str(kpi.get("frequency","")), xl + 0.1, 5.97, 2.9, 0.28, 7.5, False, _MUTED, PP_ALIGN.LEFT)

    # ── Slide 6: Marketing Strategy ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "Marketing Strategy", "Traditional & Digital Channels")
    _rect(sl, 0.3, 1.1, 6.2, 0.4, _NAVY)
    _tb(sl, "Traditional Marketing", 0.4, 1.14, 6.0, 0.3, 10, True, _WHITE, PP_ALIGN.LEFT)
    for i, t in enumerate((trad.get("strategies", []) or [])[:3]):
        yt = 1.55 + i * 1.0
        _rect(sl, 0.3, yt, 6.2, 0.93, _WHITE if i % 2 == 0 else _LIGHT)
        _tb(sl, str(t.get("tactic",""))[:52], 0.4, yt + 0.04, 4.8, 0.3, 9, True, _ORANGE, PP_ALIGN.LEFT)
        _tb(sl, str(t.get("description",""))[:105], 0.4, yt + 0.36, 4.8, 0.42, 8, False, _DARK, PP_ALIGN.LEFT)
        _tb(sl, f"₱{t.get('monthly_budget_php',0):,.0f}/mo", 5.0, yt + 0.04, 1.4, 0.3, 8, False, _AMBER, PP_ALIGN.RIGHT)
    _tb(sl, "Local Partnerships", 0.3, 4.62, 6, 0.32, 10, True, _NAVY, PP_ALIGN.LEFT)
    _bullets(sl, trad.get("local_partnerships", []), 0.3, 4.95, 6.2, 2.4, 8.5, _DARK, 4)
    _rect(sl, 6.85, 1.1, 6.15, 0.4, _ORANGE)
    _tb(sl, "Digital Marketing", 6.95, 1.14, 5.95, 0.3, 10, True, _WHITE, PP_ALIGN.LEFT)
    for i, p in enumerate((dig.get("platforms", []) or [])[:3]):
        yt = 1.55 + i * 1.0
        _rect(sl, 6.85, yt, 6.15, 0.93, _WHITE if i % 2 == 0 else _LIGHT)
        _tb(sl, str(p.get("platform",""))[:45], 6.95, yt + 0.04, 4.8, 0.3, 9, True, _ORANGE, PP_ALIGN.LEFT)
        _tb(sl, str(p.get("strategy",""))[:100], 6.95, yt + 0.36, 4.8, 0.42, 8, False, _DARK, PP_ALIGN.LEFT)
        _tb(sl, f"₱{p.get('monthly_budget_php',0):,.0f}/mo", 11.6, yt + 0.04, 1.3, 0.3, 8, False, _AMBER, PP_ALIGN.RIGHT)
    kws = dig.get("seo_keywords", [])[:5]
    _tb(sl, "SEO Keywords", 6.85, 4.62, 6, 0.32, 10, True, _NAVY, PP_ALIGN.LEFT)
    _tb(sl, " • ".join(str(k) for k in kws), 6.85, 4.95, 6.15, 0.38, 9, False, _DARK, PP_ALIGN.LEFT)
    _tb(sl, "Online Marketplaces", 6.85, 5.45, 6, 0.32, 10, True, _NAVY, PP_ALIGN.LEFT)
    _bullets(sl, dig.get("online_marketplaces", []), 6.85, 5.78, 6.15, 1.6, 9, _DARK, 3)

    # ── Slide 7: Implementation Timeline ────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "Implementation Timeline")
    phases = (tl or [])[:5]
    ph_colors = [_ORANGE, _NAVY, "#1D4ED8", "#0F766E", _AMBER]
    if phases:
        cw = 12.5 / len(phases)
        for i, ph in enumerate(phases):
            xl = 0.3 + i * cw
            bg = ph_colors[i % len(ph_colors)]
            _rect(sl, xl, 1.1, cw - 0.12, 0.55, bg)
            _tb(sl, f"Phase {ph.get('phase', i+1)}", xl + 0.05, 1.15, cw - 0.22, 0.25, 9, True, _WHITE, PP_ALIGN.CENTER)
            _tb(sl, str(ph.get("duration","")), xl + 0.05, 1.4, cw - 0.22, 0.2, 7.5, False, _WHITE, PP_ALIGN.CENTER)
            _rect(sl, xl, 1.7, cw - 0.12, 1.82, _WHITE)
            _tb(sl, str(ph.get("title",""))[:40], xl + 0.05, 1.74, cw - 0.22, 0.36, 9, True, _DARK, PP_ALIGN.LEFT)
            _bullets(sl, ph.get("key_tasks", []), xl + 0.05, 2.12, cw - 0.22, 1.3, 7.5, _DARK, 3)
            _rect(sl, xl, 3.58, cw - 0.12, 0.34, _LIGHT)
            _tb(sl, f"₱{ph.get('budget_php',0):,.0f}", xl + 0.05, 3.62, cw - 0.22, 0.26, 9, True, _ORANGE, PP_ALIGN.RIGHT)

    # ── Slide 8: Permits & Legal ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    sub8 = (f"Est. Cost: ₱{perm.get('total_estimated_cost_php_low',0):,.0f}–"
            f"₱{perm.get('total_estimated_cost_php_high',0):,.0f}  |  "
            f"{perm.get('total_processing_weeks','N/A')} processing")
    _header(sl, "Permits & Legal Requirements", sub8)
    steps8 = (perm.get("steps", []) or [])[:6]
    st_colors = [_ORANGE, _NAVY, "#1D4ED8", _GREEN, _AMBER, "#7C3AED"]
    for i, s in enumerate(steps8):
        row, col = divmod(i, 3)
        xl = 0.3 + col * 4.35
        yt = 1.1 + row * 2.85
        c = st_colors[i % len(st_colors)]
        _rect(sl, xl, yt, 0.42, 0.42, c)
        _tb(sl, str(s.get("step", i+1)), xl + 0.03, yt + 0.04, 0.36, 0.34, 12, True, _WHITE, PP_ALIGN.CENTER)
        _rect(sl, xl + 0.48, yt, 3.8, 2.68, _WHITE)
        _tb(sl, str(s.get("requirement",""))[:48], xl + 0.58, yt + 0.05, 3.6, 0.36, 9.5, True, c, PP_ALIGN.LEFT)
        _tb(sl, str(s.get("office","")), xl + 0.58, yt + 0.42, 3.6, 0.26, 7.5, False, _MUTED, PP_ALIGN.LEFT)
        cost8 = f"₱{s.get('cost_php_low',0):,.0f}–₱{s.get('cost_php_high',0):,.0f}"
        _tb(sl, cost8, xl + 0.58, yt + 0.68, 3.6, 0.26, 8.5, True, _ORANGE, PP_ALIGN.LEFT)
        _tb(sl, str(s.get("tips",""))[:95], xl + 0.58, yt + 0.95, 3.6, 1.6, 7.5, False, _DARK, PP_ALIGN.LEFT)

    # ── Slide 9: Risk Analysis ───────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    _header(sl, "Risk Analysis & Mitigation")
    risks = (risk_d.get("top_risks", []) or [])[:5]
    p_clr = {"High": _RED, "Medium": _AMBER, "Low": _GREEN}
    _rect(sl, 0.3, 1.1, 12.7, 0.42, _NAVY)
    for lbl9, cx9, cw9 in [("Risk", 0.4, 4.8), ("Probability", 5.4, 1.8), ("Impact", 7.35, 1.4), ("Mitigation Strategy", 8.9, 4.1)]:
        _tb(sl, lbl9, cx9, 1.14, cw9, 0.3, 9, True, _WHITE, PP_ALIGN.LEFT)
    for ri, r in enumerate(risks):
        yt = 1.55 + ri * 1.06
        _rect(sl, 0.3, yt, 12.7, 0.98, _WHITE if ri % 2 == 0 else _LIGHT)
        _tb(sl, str(r.get("risk",""))[:85], 0.4, yt + 0.07, 4.8, 0.84, 8.5, False, _DARK, PP_ALIGN.LEFT)
        for col9, key9, cx9, cw9 in [("probability", 5.4, 1.72), ("impact", 7.35, 1.42)]:
            val9 = r.get(key9, "Medium")
            _rect(sl, cx9, yt + 0.24, cw9, 0.4, p_clr.get(val9, _MUTED))
            _tb(sl, val9, cx9, yt + 0.28, cw9, 0.3, 8, True, _WHITE, PP_ALIGN.CENTER)
        _tb(sl, str(r.get("mitigation",""))[:130], 8.9, yt + 0.06, 4.1, 0.84, 7.5, False, _DARK, PP_ALIGN.LEFT)
    cont = risk_d.get("contingency_plans", [])[:2]
    _rect(sl, 0.3, 6.76, 12.7, 0.6, _NAVY)
    _tb(sl, "Contingency: " + "  |  ".join(str(c) for c in cont)[:145], 0.4, 6.81, 12.5, 0.45, 8, False, _WHITE, PP_ALIGN.LEFT)

    # ── Slide 10: Materials & Equipment ─────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _LIGHT)
    sub10 = (f"Equipment: ₱{mats.get('total_equipment_cost_php_low',0):,.0f}–"
             f"₱{mats.get('total_equipment_cost_php_high',0):,.0f}  |  "
             f"Monthly Materials: ₱{mats.get('total_monthly_materials_php',0):,.0f}")
    _header(sl, "Materials & Equipment", sub10)
    equip = (mats.get("equipment", []) or [])[:6]
    pr_clr = {"Essential": _RED, "Important": _AMBER, "Nice-to-have": _GREEN}
    _rect(sl, 0.3, 1.1, 8.2, 0.42, _NAVY)
    for lbl10, cx10, cw10 in [("Equipment", 0.4, 4.5), ("Cost Range", 5.0, 2.1), ("Priority", 7.2, 1.2)]:
        _tb(sl, lbl10, cx10, 1.14, cw10, 0.3, 9, True, _WHITE, PP_ALIGN.CENTER)
    for i, eq in enumerate(equip):
        yt = 1.55 + i * 0.83
        _rect(sl, 0.3, yt, 8.2, 0.77, _WHITE if i % 2 == 0 else _LIGHT)
        _tb(sl, str(eq.get("name",""))[:48], 0.4, yt + 0.05, 4.5, 0.3, 9, True, _DARK, PP_ALIGN.LEFT)
        _tb(sl, str(eq.get("where_to_buy",""))[:52], 0.4, yt + 0.42, 4.5, 0.28, 7.5, False, _MUTED, PP_ALIGN.LEFT)
        _tb(sl, f"₱{eq.get('cost_php_low',0):,.0f}–₱{eq.get('cost_php_high',0):,.0f}",
            5.0, yt + 0.2, 2.1, 0.3, 8.5, False, _ORANGE, PP_ALIGN.CENTER)
        pr = eq.get("priority", "Important")
        _rect(sl, 7.2, yt + 0.16, 1.22, 0.4, pr_clr.get(pr, _MUTED))
        _tb(sl, pr[:9], 7.2, yt + 0.2, 1.22, 0.3, 7, True, _WHITE, PP_ALIGN.CENTER)
    _rect(sl, 8.85, 1.1, 4.15, 0.42, _ORANGE)
    _tb(sl, "Monthly Materials", 8.95, 1.14, 4.0, 0.3, 9, True, _WHITE, PP_ALIGN.LEFT)
    for i, rm in enumerate((mats.get("raw_materials", []) or [])[:5]):
        yt = 1.55 + i * 0.75
        _rect(sl, 8.85, yt, 4.15, 0.68, _WHITE if i % 2 == 0 else _LIGHT)
        _tb(sl, str(rm.get("name",""))[:36], 8.95, yt + 0.04, 2.5, 0.28, 8.5, True, _DARK, PP_ALIGN.LEFT)
        _tb(sl, f"₱{rm.get('monthly_cost_php',0):,.0f}/mo", 11.35, yt + 0.04, 1.55, 0.28, 8.5, False, _ORANGE, PP_ALIGN.RIGHT)
        _tb(sl, f"{rm.get('monthly_quantity',0)} {rm.get('unit','')}", 8.95, yt + 0.36, 3.0, 0.25, 7.5, False, _MUTED, PP_ALIGN.LEFT)
    _tb(sl, "Sourcing Tips", 8.85, 5.48, 4.0, 0.3, 10, True, _NAVY, PP_ALIGN.LEFT)
    _bullets(sl, mats.get("sourcing_tips", []), 8.85, 5.78, 4.15, 1.6, 8.5, _DARK, 3)

    # ── Slide 11: Closing ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, 13.33, 7.5, _NAVY)
    _rect(sl, 0, 0, 13.33, 0.12, _ORANGE)
    _rect(sl, 0, 6.8, 13.33, 0.7, _ORANGE)
    _rect(sl, 5.75, 1.0, 0.08, 5.5, _AMBER)
    _tb(sl, "You're Ready to Launch!", 0.5, 1.2, 5.1, 0.72, 28, True, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "Your personalized blueprint covers everything you\nneed to plan, launch, and grow your business.",
        0.5, 2.0, 5.1, 1.0, 10, False, "#B0BFD8", PP_ALIGN.LEFT)
    checklist11 = [
        "✓  Business Plan & Strategy",
        "✓  Financial Projections",
        "✓  Marketing Strategy (Trad + Digital)",
        "✓  Operations Manual & Staffing",
        "✓  Permits & Legal Guide",
        "✓  Risk Management Plan",
    ]
    _bullets(sl, checklist11, 0.5, 3.1, 5.1, 3.2, 10, _AMBER, 6)
    _rect(sl, 6.15, 1.0, 6.85, 2.55, "#1A3060")
    _tb(sl, "Next Steps", 6.45, 1.1, 6.4, 0.4, 13, True, _AMBER, PP_ALIGN.LEFT)
    _bullets(sl, [
        "1. Read your full blueprint carefully",
        "2. Register your business (DTI → Barangay → BIR)",
        "3. Set your Week 1 milestone and act",
        "4. Open a GCash Business or bank account",
        "5. Contact us on WhatsApp for support",
    ], 6.45, 1.55, 6.4, 1.85, 9, _WHITE, 5)
    _rect(sl, 6.15, 3.65, 6.85, 2.5, "#152B50")
    _tb(sl, "Need Help?", 6.45, 3.75, 6.4, 0.4, 13, True, _AMBER, PP_ALIGN.LEFT)
    _tb(sl, "WhatsApp Support (24/7)", 6.45, 4.22, 6.4, 0.35, 10, False, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "support@negosyoplan.com", 6.45, 4.6, 6.4, 0.35, 10, False, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "negosyoplan.com", 6.45, 4.98, 6.4, 0.35, 10, False, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "Negosyo Plan", 0.35, 6.87, 6, 0.5, 16, True, _WHITE, PP_ALIGN.LEFT)
    _tb(sl, "Your Business Blueprint to Success", 6.5, 6.92, 6.5, 0.4, 10, False, "#FFCDA0", PP_ALIGN.RIGHT)

    filename = f"NegosyoPlan_Blueprint_{uuid.uuid4().hex[:12]}.pptx"
    prs.save(str(PRESENTATIONS_DIR / filename))
    return filename


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/")
async def root():
    return {
        "service": "Negosyo Plan AI API",
        "version": "1.0.0",
        "endpoints": [
            "POST /api/generate-blueprint",
            "POST /api/quick-analysis",
            "GET  /health",
        ],
    }


@app.get("/health")
async def health():
    return {
        "status": "healthy",
        "deepseek_configured": bool(DEEPSEEK_API_KEY),
        "model": DEEPSEEK_MODEL,
    }


@app.post("/api/generate-blueprint")
async def generate_blueprint(req: BlueprintRequest):
    """
    Full comprehensive blueprint extraction via DeepSeek.

    Accepts the business-form customization payload and returns a rich JSON
    object covering SWOT, FIFO, Risk, Financial, Strategy, Marketing,
    Timeline, Suppliers, and Legal requirements — depth scales with bundle tier.
    """
    data = await call_deepseek(
        system=SYSTEM_PROMPT,
        user=build_blueprint_prompt(req),
        max_tokens=8000,
        temperature=0.65,
        timeout=120.0,
    )
    return {
        "success": True,
        "meta": {
            "business_type": req.business_type,
            "business_type_label": req.business_type_label,
            "bundle_tier": req.bundle_tier,
            "capital_php": req.capital,
            "location": req.location.model_dump(),
        },
        "blueprint": data,
    }


@app.post("/api/quick-analysis")
async def quick_analysis(req: QuickAnalysisRequest):
    """
    Lightweight pre-purchase feasibility check (~3-5 s).

    Returns a feasibility score, revenue estimate, bundle recommendation,
    capital sufficiency note, and quick tips for the chosen business type.
    """
    data = await call_deepseek(
        system="You are a Filipino MSME advisor. Respond with valid JSON only.",
        user=build_quick_prompt(req),
        max_tokens=900,
        temperature=0.6,
        timeout=30.0,
    )
    return {"success": True, "analysis": data}


# ═══════════════════════════════════════════════════════════════════════════════
# ADMIN ROUTES
# ═══════════════════════════════════════════════════════════════════════════════

class AdminLoginRequest(BaseModel):
    username: str
    password: str


@app.post("/admin/login")
async def admin_login(creds: AdminLoginRequest):
    """Validate admin credentials and return a JWT token."""
    ok = (
        hmac.compare_digest(creds.username, ADMIN_USERNAME)
        and hmac.compare_digest(creds.password, ADMIN_PASSWORD)
    )
    if not ok:
        raise HTTPException(status_code=401, detail="Invalid username or password")
    token = create_token(creds.username)
    return {
        "success": True,
        "access_token": token,
        "username": creds.username,
        "expires_in_hours": JWT_EXPIRE_HOURS,
    }


@app.get("/admin/verify")
async def admin_verify(authorization: str | None = Header(default=None)):
    """Check whether a JWT token is still valid."""
    user = require_admin(authorization)
    return {"valid": True, "username": user}


# ── Site settings (public read, admin write) ──────────────────────────────────

@app.get("/api/settings")
async def get_settings():
    """
    Public endpoint — returns current site settings so the frontend can
    apply admin-configured prices, theme colours, promo banner, etc.
    """
    return load_settings()


@app.put("/admin/settings")
async def update_settings(
    payload: dict,
    authorization: str | None = Header(default=None),
):
    """
    Merge-update site settings.  Requires a valid admin JWT.
    Send only the keys you want to change; unchanged keys are preserved.
    """
    require_admin(authorization)
    current = load_settings()
    merged = deep_merge(current, payload)
    save_settings(merged)
    return {"success": True, "settings": merged}


@app.post("/admin/settings/reset")
async def reset_settings(authorization: str | None = Header(default=None)):
    """Restore all settings to factory defaults."""
    require_admin(authorization)
    save_settings(json.loads(json.dumps(DEFAULT_SETTINGS)))
    return {"success": True, "settings": DEFAULT_SETTINGS}


@app.post("/admin/upload-logo")
async def upload_logo(
    file: UploadFile = File(...),
    authorization: str | None = Header(default=None),
):
    """Upload a new logo image. Saves to assets/images/logo-custom.<ext>."""
    require_admin(authorization)

    allowed = {".png", ".jpg", ".jpeg", ".svg", ".webp", ".ico"}
    suffix = Path(file.filename or "logo.png").suffix.lower()
    if suffix not in allowed:
        raise HTTPException(
            400, f"File type not allowed. Use: {', '.join(sorted(allowed))}"
        )

    contents = await file.read()
    if len(contents) > 5 * 1024 * 1024:
        raise HTTPException(400, "File too large. Maximum 5 MB.")

    ASSETS_IMAGES.mkdir(parents=True, exist_ok=True)
    save_path = ASSETS_IMAGES / f"logo-custom{suffix}"
    save_path.write_bytes(contents)

    relative_url = f"assets/images/logo-custom{suffix}"
    return {
        "success": True,
        "url": relative_url,
        "filename": file.filename,
        "size_bytes": len(contents),
    }


# ── Presentation endpoints ────────────────────────────────────────────────────

@app.post("/api/generate-presentation")
async def generate_presentation(req: BlueprintRequest):
    """
    Generate a professional 11-slide PPTX from a full AI blueprint.
    Calls DeepSeek to extract blueprint data then builds the presentation.
    """
    if not PPTX_AVAILABLE:
        raise HTTPException(
            503,
            "python-pptx is not installed. Run: pip install python-pptx",
        )

    # Prune old presentations (keep last 20)
    PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)
    existing = sorted(PRESENTATIONS_DIR.glob("*.pptx"), key=lambda f: f.stat().st_mtime)
    for old in existing[:-19]:
        try:
            old.unlink()
        except Exception:
            pass

    data = await call_deepseek(
        system=SYSTEM_PROMPT,
        user=build_blueprint_prompt(req),
        max_tokens=8000,
        temperature=0.65,
        timeout=120.0,
    )

    meta = {
        "business_type": req.business_type,
        "business_type_label": req.business_type_label,
        "bundle_tier": req.bundle_tier,
        "capital_php": req.capital,
        "location": req.location.model_dump(),
    }

    try:
        filename = generate_pptx(data, meta)
    except Exception as exc:
        raise HTTPException(500, f"Presentation generation failed: {str(exc)[:300]}")

    return {
        "success": True,
        "filename": filename,
        "download_url": f"/download/presentation/{filename}",
    }


@app.get("/download/presentation/{filename}")
async def download_presentation(filename: str):
    """Serve a generated PPTX file. Rejects path traversal attempts."""
    if ".." in filename or "/" in filename or "\\" in filename:
        raise HTTPException(400, "Invalid filename")
    path = PRESENTATIONS_DIR / filename
    if not path.exists() or not path.is_file():
        raise HTTPException(404, "Presentation not found. It may have expired — please regenerate.")
    return FileResponse(
        str(path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )


# ── Genspark integration ──────────────────────────────────────────────────────

class GensparkSparkRequest(BaseModel):
    title: str
    business_type_label: str
    bundle_tier: str = "starter"
    capital: float = 50000
    location: str = "Philippines"
    description: str = ""
    outline: str


@app.post("/api/genspark-spark")
async def create_genspark_spark(req: GensparkSparkRequest):
    """
    Submit a blueprint outline to Genspark AI and return the generated
    sparkpage URL for use as the final presentation design.
    """
    if not GENSPARK_API_KEY:
        raise HTTPException(503, "Genspark API key not configured. Add GENSPARK_API_KEY to backend/.env")

    query = (
        f"Create a professional 11-slide business presentation for a {req.business_type_label} "
        f"business in {req.location} with starting capital of PHP {req.capital:,.0f}. "
        f"Bundle tier: {req.bundle_tier.upper()}. "
        f"{('Additional details: ' + req.description) if req.description else ''}\n\n"
        f"{req.outline}"
    )

    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(
                GENSPARK_URL,
                headers={
                    "Authorization": f"Bearer {GENSPARK_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "query": query,
                    "type": "presentation",
                },
            )

            if resp.status_code == 404:
                # Fallback: try alternate Genspark endpoint
                resp = await client.post(
                    "https://api.genspark.ai/v1/search",
                    headers={
                        "Authorization": f"Bearer {GENSPARK_API_KEY}",
                        "Content-Type": "application/json",
                    },
                    json={"query": query, "search_type": "presentation"},
                )

            if not resp.is_success:
                raise HTTPException(
                    resp.status_code,
                    f"Genspark API error ({resp.status_code}): {resp.text[:300]}"
                )

            data = resp.json()
            spark_url = (
                data.get("url")
                or data.get("page_url")
                or data.get("spark_url")
                or data.get("data", {}).get("url")
                or data.get("result", {}).get("url")
            )

            if not spark_url:
                raise HTTPException(500, f"Genspark returned no URL. Response: {str(data)[:300]}")

            return {
                "success": True,
                "spark_url": spark_url,
                "title": req.title,
            }

    except httpx.TimeoutException:
        raise HTTPException(504, "Genspark request timed out (120 s). Please retry.")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(500, f"Genspark integration error: {str(exc)[:300]}")
